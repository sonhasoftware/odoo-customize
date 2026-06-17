# -*- coding: utf-8 -*-

from odoo import api, fields, models
import json
import logging
import re
import json
import base64
import PyPDF2
from docx import Document
import pandas as pd
import requests
import pyodbc
import numpy as np
import io
import os

_logger = logging.getLogger(__name__)

class SonhaAIBot(models.AbstractModel):
    _name = 'sonha.ai.bot'
    _description = 'Sonha AI Bot Logic'

    @api.model
    def process_chat(self, prompt, model_name=None, record_id=None, chatter_messages=None):
        ICPSudo = self.env['ir.config_parameter'].sudo()
        provider = ICPSudo.get_param('sonha_ai.provider', 'openai')
        
        context_str = ""
        if model_name and record_id:
            try:
                record = self.env[model_name].browse(int(record_id))
                if record.exists():
                    context_str += f"Context: User is currently viewing record ID {record_id} of model '{model_name}'. "
                    if hasattr(record, 'display_name'):
                        context_str += f"Record Name: {record.display_name}. "
                    
                    if hasattr(record, 'message_ids'):
                        messages = record.message_ids.filtered(lambda m: m.message_type in ['comment', 'email'])[:10]
                        if messages:
                            context_str += "\nRecent chatter messages for context:\n"
                            for msg in messages:
                                body = re.sub('<[^<]+>', '', msg.body) if msg.body else ''
                                context_str += f"[{msg.date}] {msg.author_id.name or 'System'}: {body}\n"
            except Exception as e:
                _logger.warning("Sonha AI: Could not fetch record context: %s", str(e))
                
        if chatter_messages:
            clean_chatter = []
            for msg in chatter_messages:
                # Remove json blocks from history to prevent in-context learning poisoning
                if not (('{' in msg and '}' in msg) or '```json' in msg):
                    clean_chatter.append(msg)
            if clean_chatter:
                context_str += "\nRecent chatter messages for context:\n" + "\n".join(clean_chatter)

        # Luôn sử dụng Vector DB (SQL Server) cho mọi mô hình AI (Gemini, OpenAI, Ollama)
        vector_context = self._search_sql_vectors(prompt)
        if vector_context:
            context_str += "\n\n[KIẾN THỨC TỪ TÀI LIỆU NỘI BỘ CỦA CÔNG TY]:\n" + vector_context

        
        if provider == 'ollama':
            system_prompt = """Bạn là một trợ lý AI thông minh được tích hợp bên trong hệ thống Odoo ERP.
Nhiệm vụ chính của bạn là đọc các dữ liệu từ Cơ sở tri thức (Knowledge Base) do hệ thống cung cấp và trả lời câu hỏi của người dùng một cách chính xác.
CÁC QUY TẮC BẮT BUỘC (TUYỆT ĐỐI TUÂN THỦ):
1. Bạn BẮT BUỘC phải luôn trả lời bằng Tiếng Việt (Vietnamese). Không bao giờ sử dụng tiếng Anh, tiếng Nga hay ngôn ngữ khác.
2. KHÔNG BAO GIỜ được trả về định dạng JSON hoặc code. Chỉ trả lời bằng các đoạn văn bản ngôn ngữ tự nhiên bình thường.
3. Nếu người dùng yêu cầu tóm tắt hoặc trích xuất thông tin, hãy viết một đoạn văn ngắn gọn, dễ hiểu.
4. Trả lời lịch sự, chuyên nghiệp và đi thẳng vào trọng tâm.
"""
        else:
            system_prompt = """You are a helpful AI assistant embedded inside an Odoo ERP system.
If the user's request is a command to open, navigate to, or view a specific screen, list, or menu (e.g., 'mở danh sách nhân viên', 'đi tới kho', 'mở đặt xe'), you MUST return your response starting exactly with the text '[NAVIGATE]: ' followed by the exact short keyword of that screen AS TYPED BY THE USER. Do not translate or change the user's words into standard Odoo module names. If the user says "đặt xe", return "[NAVIGATE]: đặt xe". Example: [NAVIGATE]: Bán hàng

If the user asks for a chart, a report, or to export data, you MUST act as an Odoo Data Analyst. Generate a pure JSON block providing Odoo ORM parameters (model, domain, fields, groupby) instead of SQL.
CRITICAL SCHEMA RULE: When querying, strictly use Odoo model names, NOT table names:
- Khách hàng / Liên hệ / Đối tác -> Model: res.partner
- Nhân viên -> Model: hr.employee
- Đơn bán hàng -> Model: sale.order
- Đơn mua hàng -> Model: purchase.order
- Sản phẩm -> Model: product.template or product.product
Do NOT invent fields that are not standard in Odoo.
YOU MUST ALWAYS RESPOND ONLY WITH THE JSON BLOCK. DO NOT EXPLAIN.

For a chart, respond ONLY with this JSON block format:
{
  "action": "orm_chart",
  "model": "hr.employee",
  "domain": [],
  "groupby": ["department_id"],
  "chart_type": "horizontalBar"
}
Note: chart_type can be 'bar', 'horizontalBar', 'line', 'pie', or 'doughnut'. If the user asks to group by multiple fields (max 2), provide them in the groupby list (e.g. ["company_id", "department_id"]).

For an export, respond ONLY with this JSON block format:
{
  "action": "orm_export",
  "model": "hr.employee",
  "domain": [],
  "fields": ["name", "work_email", "department_id"],
  "filename": "nhan_vien.xlsx"
}

If the user asks to create a new record OR asks how to add/create one (e.g. "Tạo khách hàng mới tên ABC", "Tạo đơn bán hàng", "thêm mới nhân viên như thế nào"), respond ONLY with this JSON block format to provide a pre-filled demo form:
{
  "action": "orm_create_draft",
  "model": "res.partner",
  "context": {
    "default_name": "Nguyễn Văn A (Demo)",
    "default_phone": "0987654321"
  }
}
Note: Extract entities from user prompt or invent reasonable DEMO data if not provided. Put them as 'default_[fieldname]' in the context dictionary. Do NOT invent fields. (e.g. hr.employee fields: default_name, default_work_phone, default_work_email).

--- TÍNH NĂNG TỐI ƯU VĂN BẢN ---
Nếu người dùng bắt đầu câu hỏi bằng một trong các thẻ lệnh: [Tối ưu], [Trang trọng], [Thân thiện], [Tóm tắt], bạn phải đóng vai là chuyên gia tối ưu hóa văn bản nội bộ và soạn thảo email thương mại cho Odoo.
- [Tối ưu]: Hãy viết lại đoạn văn sao cho chuyên nghiệp, mạch lạc hơn, sửa toàn bộ lỗi chính tả và ngữ pháp.
- [Trang trọng]: Chuyển đổi giọng điệu sang lịch sự, phù hợp gửi cho đối tác/khách hàng/cấp trên.
- [Thân thiện]: Chuyển đổi giọng điệu sang cởi mở, dùng cho trao đổi nội bộ đội ngũ.
- [Tóm tắt]: Rút gọn văn bản thành một đoạn ngắn không quá 3 câu nhưng giữ đủ ý chính.
Yêu cầu đầu ra cho tính năng này: 
- Chỉ trả về DUY NHẤT đoạn văn bản đã được xử lý (hoặc mã HTML sạch nếu văn bản gốc có format).
- Không giải thích, không kèm theo lời thoại của AI như "Đây là đoạn văn của bạn:".

--- TÍNH NĂNG PHÂN LOẠI TÀI LIỆU ---
Nếu người dùng CỐ TÌNH gõ thẻ lệnh [Phân loại] hoặc [Extract] ở đầu câu hỏi, bạn phải đóng vai là bộ định tuyến dữ liệu tự động (Document Classifier) của Odoo. 
Nhiệm vụ của bạn là phân tích đoạn văn bản được trích xuất từ một tài liệu tải lên (Hóa đơn, Hợp đồng, CV, Báo giá...) và trả về định dạng JSON để hệ thống tự động xử lý.
Dựa trên nội dung tài liệu, hãy xác định các thông tin sau:
1. document_type: Loại tài liệu (Chọn 1 trong các giá trị: 'invoice', 'contract', 'recruitment', 'other').
2. partner_name: Tên của đối tác/khách hàng/ứng viên xuất hiện trong tài liệu (nếu có).
3. total_amount: Tổng số tiền nếu là hóa đơn/báo giá (nếu không có điền null).
4. summary: Tóm tắt nội dung tài liệu trong vòng 1 câu.
Yêu cầu đầu ra bắt buộc: 
- Chỉ trả về một chuỗi JSON hợp lệ, không bọc trong block code ```json...```. Không trả thêm bất kỳ ký tự nào khác.
Ví dụ đầu ra: {"document_type": "invoice", "partner_name": "Công ty ABC", "total_amount": 15000000, "summary": "Hóa đơn thanh toán dịch vụ phần mềm tháng 5"}

For all other general questions, answer normally in natural language.
CRITICAL INSTRUCTION FOR AI: 
- DO NOT OUTPUT JSON unless the user explicitly triggers one of the special commands above ([NAVIGATE], [Extract], etc.).
- If the user asks a normal question about the Knowledge Base, YOU MUST ANSWER IN PLAIN NATURAL LANGUAGE (Vietnamese). DO NOT USE JSON."""
        full_prompt = f"{system_prompt}\n\n{context_str}\n\nUser Question: {prompt}"

        try:
            if provider == 'ollama':
                resp = self._call_ollama_chat(full_prompt)
            elif provider == 'openai':
                resp = self._call_openai(full_prompt)
            elif provider == 'gemini':
                resp = self._call_gemini(full_prompt)
            else:
                resp = {"error": "Invalid AI provider configured."}
            
            # Map error to result so UI can display it
            if 'error' in resp and 'result' not in resp:
                resp['result'] = f"Lỗi từ hệ thống AI ({provider}): {resp['error']}"
            
            if 'result' in resp and '[NAVIGATE]:' in resp['result']:
                keyword = resp['result'].split('[NAVIGATE]:')[1].strip()
                menu = self.env['ir.ui.menu'].sudo().search([('name', 'ilike', keyword)], limit=1)
                
                if not menu:
                    action = self.env['ir.actions.act_window'].sudo().search([('name', 'ilike', keyword)], limit=1)
                    if action:
                        menu = self.env['ir.ui.menu'].sudo().search([('action', '=', f'ir.actions.act_window,{action.id}')], limit=1)
                        if not menu:
                            resp['result'] = f"Đang chuyển hướng đến: {action.name}..."
                            resp['action_id'] = action.id
                
                if not menu and 'action_id' not in resp:
                    words = [w for w in keyword.split() if len(w) > 2]
                    if words:
                        domain = []
                        for w in words:
                            domain.append(('name', 'ilike', w))
                        menu = self.env['ir.ui.menu'].sudo().search(domain, limit=1)
                        
                if not menu and 'action_id' not in resp:
                    words = [w for w in keyword.split() if len(w) > 2]
                    for w in reversed(words):
                        menu = self.env['ir.ui.menu'].sudo().search([('name', 'ilike', w)], limit=1)
                        if menu:
                            break
                            
                if menu:
                    resp['result'] = f"Đang chuyển hướng đến: {menu.name}..."
                    if menu.action:
                        resp['action_id'] = menu.action.id
                    resp['menu_id'] = menu.id
                elif 'action_id' not in resp:
                    resp['result'] = f"Tôi không tìm thấy chức năng nào tên là '{keyword}' trong hệ thống."
            
            if 'result' in resp:
                text = resp['result']
                json_str = None
                
                # First try to find markdown json block
                json_match = re.search(r'```(?:json)?\s*(\{.*?\})\s*```', text, re.DOTALL | re.IGNORECASE)
                if json_match:
                    json_str = json_match.group(1)
                else:
                    # Fallback to finding first { and last }
                    start = text.find('{')
                    end = text.rfind('}')
                    if start != -1 and end != -1 and end > start:
                        json_str = text[start:end+1]
                        
                if json_str:
                    try:
                        ai_json = json.loads(json_str)
                        if isinstance(ai_json, dict):
                            action_type = ai_json.get('action')
                            if action_type == 'orm_create_draft':
                                resp['action'] = 'orm_create_draft'
                                resp['model'] = ai_json.get('model')
                                resp['context'] = ai_json.get('context', {})
                                resp['result'] = "Đang chuẩn bị form dữ liệu nháp. Vui lòng kiểm tra và nhấn Lưu..."
                                return resp
                                
                            if action_type in ['orm_chart', 'orm_export']:
                                model_name = ai_json.get('model')
                                if not model_name or model_name not in self.env:
                                    resp['result'] = f"Lỗi: Model {model_name} không tồn tại hoặc AI không xác định được."
                                    return resp
                                    
                                try:
                                    model_obj = self.env[model_name]
                                    domain = ai_json.get('domain', [])
                                
                                    data = []
                                    keys = []
                                    
                                    if ai_json['action'] == 'orm_chart':
                                        groupby = ai_json.get('groupby', [])
                                        if not groupby:
                                            resp['result'] = "Lỗi: Biểu đồ yêu cầu ít nhất một trường để nhóm (groupby)."
                                            return resp
                                            
                                        # Sử dụng read_group để lấy dữ liệu thống kê an toàn qua ORM
                                        read_group_res = model_obj.read_group(
                                            domain=domain,
                                            fields=[],
                                            groupby=groupby,
                                            lazy=True
                                        )
                                        
                                        # Định dạng lại kết quả cho logic vẽ chart cũ
                                        for res in read_group_res:
                                            row = {}
                                            for gb in groupby:
                                                val = res.get(gb)
                                                if isinstance(val, tuple) and len(val) == 2:
                                                    row[gb] = val[1]
                                                elif val and isinstance(val, int) and gb in model_obj._fields and model_obj._fields[gb].type == 'many2one':
                                                    try:
                                                        rec = self.env[model_obj._fields[gb].comodel_name].sudo().browse(val)
                                                        row[gb] = rec.display_name or str(val)
                                                    except:
                                                        row[gb] = str(val)
                                                elif gb in model_obj._fields and model_obj._fields[gb].type == 'selection':
                                                    selection_dict = dict(model_obj._fields[gb]._description_selection(self.env))
                                                    row[gb] = selection_dict.get(val, val) if val else val
                                                elif gb in model_obj._fields and model_obj._fields[gb].type == 'boolean':
                                                    row[gb] = "Có" if val else "Không"
                                                else:
                                                    row[gb] = val
                                            
                                            # Tìm trường count phù hợp (trong Odoo, read_group lazy=True trả về "{groupby}_count")
                                            count_val = res.get('__count', 0)
                                            if not count_val and groupby:
                                                count_val = res.get(f"{groupby[0]}_count", 0)
                                            row['total_count'] = count_val
                                            data.append(row)
                                            
                                        if groupby:
                                            # Khóa keys giống format cũ: [groupby1, groupby2, total_count]
                                            keys = groupby + ['total_count']
    
                                    elif ai_json['action'] == 'orm_export':
                                        fields_to_read = ai_json.get('fields', [])
                                        search_res = model_obj.search_read(
                                            domain=domain,
                                            fields=fields_to_read,
                                            limit=1000
                                        )
                                        # Format data
                                        for res in search_res:
                                            row = {}
                                            for f in fields_to_read:
                                                val = res.get(f)
                                                if isinstance(val, tuple) and len(val) == 2:
                                                    row[f] = val[1]
                                                else:
                                                    row[f] = val
                                            data.append(row)
                                        keys = fields_to_read
    
                                    if ai_json['action'] == 'orm_chart':
                                        if data and len(keys) >= 2:
                                            def parse_label(val):
                                                if val is None: return "None"
                                                s = str(val)
                                                if s.startswith('{') and s.endswith('}'):
                                                    try:
                                                        import ast
                                                        d = ast.literal_eval(s)
                                                        if isinstance(d, dict):
                                                            return d.get('vi_VN', d.get('en_US', list(d.values())[0] if d else s))
                                                    except Exception:
                                                        pass
                                                return s
                                                
                                            def parse_value(val):
                                                import decimal
                                                if isinstance(val, decimal.Decimal):
                                                    return float(val)
                                                if isinstance(val, (int, float)):
                                                    return val
                                                try:
                                                    return float(val)
                                                except (ValueError, TypeError):
                                                    return 0
    
                                            colors = [
                                                "#f72585", "#7209b7", "#3a0ca3", "#4361ee", "#4cc9f0",
                                                "#f94144", "#f3722c", "#f8961e", "#f9844a", "#f9c74f",
                                                "#90be6d", "#43aa8b", "#4d908e", "#577590", "#277da1"
                                            ]
                                            chart_type = ai_json.get('chart_type', 'bar')
                                            
                                            if len(keys) >= 3:
                                                from itertools import groupby
                                                data.sort(key=lambda x: str(x.get(keys[0], '')))
                                                grouped_charts = []
                                                for group_name, group_data in groupby(data, key=lambda x: str(x.get(keys[0], ''))):
                                                    group_list = list(group_data)
                                                    labels = [parse_label(row.get(keys[1])) for row in group_list]
                                                    values = [parse_value(row.get(keys[2])) for row in group_list]
                                                    
                                                    # Gom nhóm các giá trị nhỏ nếu quá nhiều
                                                    if len(labels) > 12:
                                                        items = list(zip(labels, values))
                                                        items.sort(key=lambda x: x[1] if isinstance(x[1], (int, float)) else 0, reverse=True)
                                                        top_items = items[:11]
                                                        other_value = sum(x[1] if isinstance(x[1], (int, float)) else 0 for x in items[11:])
                                                        labels = [x[0] for x in top_items] + ["Khác (Others)"]
                                                        values = [x[1] for x in top_items] + [other_value]
                                                    if chart_type in ['pie', 'doughnut', 'polarArea']:
                                                        bg_colors = [colors[i % len(colors)] for i in range(len(labels))]
                                                    else:
                                                        bg_colors = ["#d32f2f"] * len(labels)
                                                    grouped_charts.append({
                                                        "title": str(parse_label(group_name)),
                                                        "labels": labels,
                                                        "datasets": [{
                                                            "label": str(keys[2]),
                                                            "data": values,
                                                            "backgroundColor": bg_colors,
                                                            "borderWidth": 1,
                                                            "borderColor": "#ffffff"
                                                        }],
                                                        "chart_type": chart_type
                                                    })
                                                resp['result'] = "Dữ liệu biểu đồ của bạn đã được chia nhóm:"
                                                resp['multiChartData'] = grouped_charts
                                            else:
                                                labels = [parse_label(row.get(keys[0])) for row in data]
                                                values = [parse_value(row.get(keys[1])) for row in data]
                                                
                                                # Gom nhóm các giá trị nhỏ nếu quá nhiều
                                                if len(labels) > 12:
                                                    items = list(zip(labels, values))
                                                    items.sort(key=lambda x: x[1] if isinstance(x[1], (int, float)) else 0, reverse=True)
                                                    top_items = items[:11]
                                                    other_value = sum(x[1] if isinstance(x[1], (int, float)) else 0 for x in items[11:])
                                                    labels = [x[0] for x in top_items] + ["Khác (Others)"]
                                                    values = [x[1] for x in top_items] + [other_value]
                                                if chart_type in ['pie', 'doughnut', 'polarArea']:
                                                    bg_colors = [colors[i % len(colors)] for i in range(len(labels))]
                                                else:
                                                    bg_colors = ["#d32f2f"] * len(labels)
                                                    
                                                resp['result'] = "Dữ liệu biểu đồ của bạn đã sẵn sàng:"
                                                resp['chartData'] = {
                                                    "labels": labels,
                                                    "datasets": [{
                                                        "label": str(keys[1]),
                                                        "data": values,
                                                        "backgroundColor": bg_colors,
                                                        "borderWidth": 1,
                                                        "borderColor": "#ffffff"
                                                    }],
                                                    "chart_type": chart_type
                                                }
                                        else:
                                            resp['result'] = "Không có đủ dữ liệu để vẽ biểu đồ."
                                    
                                    elif ai_json['action'] == 'orm_export':
                                        try:
                                            import xlsxwriter
                                            output = io.BytesIO()
                                            workbook = xlsxwriter.Workbook(output)
                                            worksheet = workbook.add_worksheet()
                                            if data:
                                                for col, key in enumerate(keys):
                                                    worksheet.write(0, col, key)
                                                for row, record in enumerate(data):
                                                    for col, key in enumerate(keys):
                                                        worksheet.write(row + 1, col, str(record.get(key, '')))
                                            workbook.close()
                                            output.seek(0)
                                            excel_b64 = base64.b64encode(output.read()).decode()
                                            resp['result'] = "File dữ liệu của bạn đã được tạo thành công:"
                                            resp['fileData'] = excel_b64
                                            resp['filename'] = ai_json.get('filename', 'export_data.xlsx')
                                        except ImportError:
                                            resp['result'] = "Thiếu thư viện xlsxwriter trên server. Vui lòng cài đặt: pip install xlsxwriter"
                                except Exception as orm_ex:
                                    _logger.error("ORM execution error: %s", str(orm_ex))
                                    resp['result'] = f"Lỗi thực thi ORM: {str(orm_ex)}"
                    except Exception as ex:
                        _logger.error("JSON parse error: %s", str(ex))
                        pass
            return resp
            
        except Exception as e:
            _logger.error("Sonha AI General Error: %s", str(e))
            return {'result': f"Lỗi hệ thống: {str(e)}"}

    def _get_sql_connection(self):
        ICPSudo = self.env['ir.config_parameter'].sudo()
        server = ICPSudo.get_param('sonha_ai.sql_server', 'LAPTOP-BP99BUKE\\SQLEXPRESS')
        database = ICPSudo.get_param('sonha_ai.sql_database', 'OdooAI_Vector')
        user = ICPSudo.get_param('sonha_ai.sql_user', '')
        password = ICPSudo.get_param('sonha_ai.sql_password', '')

        driver_name = "SQL Server"
        for d in pyodbc.drivers():
            if "ODBC Driver" in d and "SQL Server" in d:
                driver_name = d
                break
                
        if user and password:
            conn_str = f"DRIVER={{{driver_name}}};SERVER={server};DATABASE={database};UID={user};PWD={password};"
        else:
            conn_str = f"DRIVER={{{driver_name}}};SERVER={server};DATABASE={database};Trusted_Connection=yes;"
            
        return pyodbc.connect(conn_str)

    def _generate_ollama_embedding(self, text):
        ICPSudo = self.env['ir.config_parameter'].sudo()
        url = ICPSudo.get_param('sonha_ai.ollama_url', 'http://localhost:11434')
        model = ICPSudo.get_param('sonha_ai.ollama_embed_model', 'nomic-embed-text')
        
        resp = requests.post(f"{url}/api/embeddings", json={
            "model": model,
            "prompt": text
        }, timeout=30)
        if resp.status_code == 200:
            return resp.json().get('embedding')
        return []

    def _chunk_text(self, text, chunk_size=500):
        words = text.split()
        chunks = []
        current_chunk = []
        current_length = 0
        for word in words:
            if current_length + len(word) > chunk_size and current_chunk:
                chunks.append(" ".join(current_chunk))
                current_chunk = []
                current_length = 0
            current_chunk.append(word)
            current_length += len(word) + 1
        if current_chunk:
            chunks.append(" ".join(current_chunk))
        return chunks

    def _sync_knowledge_to_sql(self):
        """Đồng bộ toàn bộ file kiến thức từ Odoo sang SQL Server"""
        company = self.env.company
        if not company.sonha_ai_knowledge_attachment_ids:
            return
            
        try:
            conn = self._get_sql_connection()
            cursor = conn.cursor()
            
            # Xóa sạch data cũ
            cursor.execute("TRUNCATE TABLE ai_document_chunks")
            
            for attachment in company.sonha_ai_knowledge_attachment_ids:
                try:
                    file_content = base64.b64decode(attachment.datas)
                    file_ext = os.path.splitext(attachment.name)[1].lower()
                    text_content = ""
                    
                    if file_ext == '.pdf':
                        pdf_reader = PyPDF2.PdfReader(io.BytesIO(file_content))
                        for page in pdf_reader.pages:
                            text_content += page.extract_text() or ""
                    elif file_ext == '.docx':
                        doc = Document(io.BytesIO(file_content))
                        text_content = "\n".join([para.text for para in doc.paragraphs])
                    elif file_ext == '.txt':
                        text_content = file_content.decode('utf-8', errors='ignore')
                    elif file_ext == '.csv':
                        text_content = file_content.decode('utf-8', errors='ignore')
                    elif file_ext in ['.xls', '.xlsx']:
                        import openpyxl
                        wb = openpyxl.load_workbook(io.BytesIO(file_content), data_only=True)
                        sheet = wb.active
                        rows = []
                        for row in sheet.iter_rows(values_only=True, max_row=500):
                            rows.append("\t".join([str(v) if v is not None else "" for v in row]))
                        text_content = "\n".join(rows)
                    elif file_ext in ['.ppt', '.pptx']:
                        import pptx
                        prs = pptx.Presentation(io.BytesIO(file_content))
                        slides_text = []
                        for slide in prs.slides:
                            for shape in slide.shapes:
                                if hasattr(shape, "text"):
                                    slides_text.append(shape.text)
                        text_content = "\n".join(slides_text)
                        
                    if text_content:
                        chunks = self._chunk_text(text_content, 500)
                        for chunk in chunks:
                            embedding = self._generate_ollama_embedding(chunk)
                            if embedding:
                                emb_str = json.dumps(embedding)
                                cursor.execute(
                                    "INSERT INTO ai_document_chunks (document_name, chunk_text, embedding) VALUES (?, ?, ?)",
                                    (attachment.name, chunk, emb_str)
                                )
                except Exception as e:
                    _logger.error(f"Error processing file {attachment.name}: {e}")
                    
            conn.commit()
            conn.close()
        except Exception as e:
            _logger.error(f"SQL Sync Error: {e}")


    def _search_sql_vectors(self, query, top_k=5):
        embedding = self._generate_ollama_embedding(query)
        if not embedding:
            return ""
            
        try:
            conn = self._get_sql_connection()
            cursor = conn.cursor()
            cursor.execute("SELECT document_name, chunk_text, embedding FROM ai_document_chunks")
            rows = cursor.fetchall()
            
            results = []
            import numpy as np
            query_vec = np.array(embedding)
            
            for row in rows:
                if row[2]:
                    import json
                    doc_vec = np.array(json.loads(row[2]))
                    # Cosine similarity
                    similarity = np.dot(query_vec, doc_vec) / (np.linalg.norm(query_vec) * np.linalg.norm(doc_vec))
                    results.append((similarity, row[0], row[1]))
            
            results.sort(key=lambda x: x[0], reverse=True)
            top_results = results[:top_k]
            
            context = ""
            for score, doc_name, text in top_results:
                if score > 0.3: # Threshold
                    context += f"[Knowledge Base: {doc_name}]\\n{text}\\n\\n"
            
            conn.close()
            return context
        except Exception as e:
            return f"(Lỗi tìm kiếm Vector: {str(e)})"

    def _call_ollama_chat(self, prompt):
        ICPSudo = self.env['ir.config_parameter'].sudo()
        url = ICPSudo.get_param('sonha_ai.ollama_url', 'http://localhost:11434')
        model = ICPSudo.get_param('sonha_ai.ollama_model', 'qwen2.5:3b')
        
        try:
            import requests
            resp = requests.post(f"{url}/api/generate", json={
                "model": model,
                "prompt": prompt,
                "stream": False
            }, timeout=120)
            
            if resp.status_code == 200:
                data = resp.json()
                return {'result': data.get('response', '')}
            else:
                return {'result': f"Lỗi từ Ollama: {resp.text}"}
        except Exception as e:
            return {'result': f"Lỗi kết nối Ollama: {str(e)}"}

    def _call_openai(self, prompt):
        try:
            import openai
        except ImportError:
            return {"error": "OpenAI library is not installed. Please pip install openai."}
            
        ICPSudo = self.env['ir.config_parameter'].sudo()
        api_key = ICPSudo.get_param('sonha_ai.openai_api_key')
        model = ICPSudo.get_param('sonha_ai.openai_model', 'gpt-4o-mini')
        
        if not api_key:
            return {"error": "OpenAI API Key is not configured in Settings."}

        try:
            client = openai.OpenAI(api_key=api_key)
            response = client.chat.completions.create(
                model=model,
                messages=[
                    {"role": "user", "content": prompt}
                ],
                max_tokens=1000,
            )
            return {"result": response.choices[0].message.content}
        except Exception as e:
            return {"error": f"OpenAI API Error: {str(e)}"}

    def _call_gemini(self, prompt):
        try:
            import google.generativeai as genai
        except ImportError:
            return {"error": "google-generativeai library is not installed. Please pip install google-generativeai."}
            
        ICPSudo = self.env['ir.config_parameter'].sudo()
        api_key = ICPSudo.get_param('sonha_ai.gemini_api_key')
        model_name = ICPSudo.get_param('sonha_ai.gemini_model', 'gemini-1.5-flash')
        
        if not api_key:
            return {"error": "Google Gemini API Key is not configured in Settings."}

        try:
            genai.configure(api_key=api_key)
            model = genai.GenerativeModel(model_name)
            response = model.generate_content(prompt)
            
            try:
                result_text = response.text
            except ValueError as ve:
                _logger.warning("Gemini text accessor failed: %s", str(ve))
                # Check if candidates exist and grab finish reason
                if response.candidates:
                    reason = response.candidates[0].finish_reason
                    # If finish reason is STOP (1) or others but no text, return empty or default
                    result_text = "*(AI đã xử lý xong nhưng không trả về bất kỳ văn bản nào. Xin vui lòng đổi lại cách hỏi hoặc kiểm tra Prompt)*"
                else:
                    return {"error": "Gemini returned an completely empty response."}
                    
            return {"result": result_text}
        except Exception as e:
            error_msg = f"Gemini API Error: {str(e)}"
            return {"error": error_msg}
